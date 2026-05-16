from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_summary_slide():
    prs = Presentation()
    
    # Use a wide aspect ratio for a modern look
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 1. Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Executive Summary: Predicting Customer Booking Behavior"
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.LEFT

    # 2. Key Metrics Box
    shape = slide.shapes.add_shape(6, Inches(0.5), Inches(1.2), Inches(3.5), Inches(1.5)) # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 32, 77) # Dark Blue
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = "Model Accuracy: 96%"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "ROC-AUC Score: 0.997"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    # 3. Text Findings (Left Column)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(5.5), Inches(4))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    def add_bullet(text, level=0, bold=False):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(16)
        p.font.bold = bold

    add_bullet("Trip Type Preferences:", bold=True)
    add_bullet("Circle Trips: Highest demand for extra baggage (78%).", level=1)
    add_bullet("Round Trips: 43% desire in-flight meals & preferred seats.", level=1)
    add_bullet("Groups (3-5): Strongest preference for seat selection.", level=1)
    
    add_bullet("Booking Insights:", bold=True)
    add_bullet("Decisive Window: Conversions peak at ~80 days lead time.", level=1)
    add_bullet("Top Markets: Malaysia, Australia, and China.", level=1)
    
    add_bullet("Strategic Recommendation:", bold=True)
    add_bullet("Target customers in the 60-80 day window with personalized bundles (Baggage + Seat).", level=1)

    # 4. Feature Importance Image (Right Column)
    slide.shapes.add_picture('booking_factors.png', Inches(6.5), Inches(1.2), width=Inches(6.3))

    # 5. Legend/Footer
    footerBox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.5))
    tf = footerBox.text_frame
    p = tf.add_paragraph()
    p.text = "Source: British Airways Summer Schedule & Booking Datasets | Random Forest Classifier (Balanced Classes)"
    p.font.size = Pt(10)
    p.font.italic = True

    prs.save('Booking_Prediction_Summary.pptx')
    print("PowerPoint slide created: Booking_Prediction_Summary.pptx")

if __name__ == "__main__":
    create_summary_slide()
